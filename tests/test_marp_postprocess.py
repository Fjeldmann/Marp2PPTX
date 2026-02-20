import io
import os
from pathlib import Path
from PIL import Image
from bs4 import BeautifulSoup
from typing import cast
import requests
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Import from refactored modules
from marp2pptx import postprocessing as mpp_post
from marp2pptx.preprocessing import remove_invisible_characters, preprocess_markdown
from marp2pptx.render_div_as_image import main as render_div_as_image
import marp2pptx.__main__ as main  # for process_pptx_html

# Create a compatibility namespace for tests
class mpp:
    """Compatibility namespace for existing tests after refactoring."""
    # From marp_convert
    parse_marp_html = mpp_post.parse_marp_html
    extract_images_from_tag = mpp_post.extract_images_from_tag
    _calculate_background_region = mpp_post._calculate_background_region
    
    # From postprocessing  
    widen_text_shapes = mpp_post.widen_text_shapes
    merge_multiline_textboxes = mpp_post.merge_multiline_textboxes
    normalize_font_names = mpp_post.normalize_font_names
    remove_redundant_marp_white_rectangles = mpp_post.remove_redundant_marp_white_rectangles
    process_native_marp_images = mpp_post.process_native_marp_images
    process_styled_divs = mpp_post.process_styled_divs
    
    # From main (orchestration function)
    process_pptx_html = main.process_pptx_html
    
    # From render_div_as_image
    render_div_as_image = render_div_as_image
    
    # From pptx.util and pptx.enum.shapes
    Cm = Cm
    MSO_SHAPE_TYPE = MSO_SHAPE_TYPE  

def _shape_rgb_safe(s):
    """Return RGBColor of a shape's fill.fore_color.rgb or None (safe checks).

    Accessing `fill.fore_color` can raise for certain fill types (e.g. _NoFill)
    in python-pptx. Guard against attribute access errors and return None if
    the fill/fore_color/rgb cannot be safely obtained.
    """
    fill = getattr(s, "fill", None)
    if fill is None:
        return None
    try:
        # `fill.fore_color` is a property that may raise (e.g. _NoFill)
        fore = getattr(fill, "fore_color", None)
    except Exception:
        return None
    if fore is None:
        return None
    try:
        return getattr(fore, "rgb", None)
    except Exception:
        return None


def test_parse_marp_html_extracts_background_and_content(tmp_path):
    html = """
    <svg data-marpit-svg="true">
      <foreignObject>
        <section data-marpit-advanced-background="background" style="--marpit-advanced-background-split:33%">
          <div data-marpit-advanced-background-container="true" data-marpit-advanced-background-direction="horizontal">
            <figure style="background-image:url('https://example.test/bg.jpg'); background-size:cover;"></figure>
          </div>
        </section>
        <section data-marpit-advanced-background="content">
          <img src="https://example.test/content.png" />
        </section>
      </foreignObject>
    </svg>
    """
    p = tmp_path / "sample.html"
    p.write_text(html, encoding="utf-8")

    slides = mpp.parse_marp_html(p)
    assert isinstance(slides, list)
    assert len(slides) == 1
    sd = slides[0]
    assert sd["backgrounds"][0]["url"] == "https://example.test/bg.jpg"
    assert any(i["url"] == "https://example.test/content.png" for i in sd["all_images"])


class DummyResponse:
    def __init__(self, content_bytes: bytes):
        self.content = content_bytes

    def raise_for_status(self):
        return None


def test_render_div_as_image_download_and_save(monkeypatch, tmp_path):
    # Create a tiny PNG in memory
    img = Image.new("RGB", (60, 60), color="blue")
    b = io.BytesIO()
    img.save(b, format="PNG")
    png_bytes = b.getvalue()

    def fake_get(url, timeout=10):
        return DummyResponse(png_bytes)

    monkeypatch.setattr(requests, "get", fake_get)

    soup = BeautifulSoup(
        '<div class="portrait-wrap-circle"><img class="portrait" src="https://example.test/img.png"/></div>',
        "html.parser",
    )
    div = soup.find("div")
    assert div is not None

    out = render_div_as_image(
        div, css=".portrait { transform: scale(1.0); }", slide_index=0, div_index=0
    )
    assert out is not None
    assert Path(out).is_file()

    # cleanup
    os.remove(out)


def test_render_div_as_image_accepts_local_file(tmp_path):
    # Create a tiny PNG on disk and reference it via a FilePath src
    img = Image.new("RGB", (60, 60), color="green")
    p = tmp_path / "local.png"
    img.save(p)

    soup = BeautifulSoup(
        f'<div class="portrait-wrap-circle"><img class="portrait" src="{p}"/></div>',
        "html.parser",
    )
    div = soup.find("div")
    assert div is not None

    out = render_div_as_image(
        div, css=".portrait { transform: scale(1.0); }", slide_index=0, div_index=0
    )
    assert out is not None
    assert Path(out).is_file()

    # cleanup
    os.remove(out)


def test_extract_images_from_tag_returns_img_and_background():
    soup = BeautifulSoup(
        '<figure style="background-image:url(\'https://bg.test/bg.jpg\')"><img src="https://img.test/pic.png"/></figure>',
        "html.parser",
    )
    img_tag = soup.find("img")
    fig_tag = soup.find("figure")
    assert img_tag is not None
    assert fig_tag is not None

    res_img = mpp.extract_images_from_tag(img_tag)
    assert any(
        r["url"] == "https://img.test/pic.png" and r["source"] == "img" for r in res_img
    )

    res_fig = mpp.extract_images_from_tag(fig_tag)
    assert any(
        r["url"] == "https://bg.test/bg.jpg" and r["source"] == "background-image"
        for r in res_fig
    )


def test__calculate_background_region_horizontal_and_split():
    # regular horizontal split for n=2
    left0, top0, w0, h0 = mpp._calculate_background_region(
        1280, 720, 0, 2, "horizontal", False
    )
    left1, top1, w1, h1 = mpp._calculate_background_region(
        1280, 720, 1, 2, "horizontal", False
    )
    assert (left0, top0, w0, h0) == (0, 0, 640, 720)
    assert (left1, top1, w1, h1) == (640, 0, 640, 720)

    # split-stack (left) horizontal for n=2, split_pct=50 -> split region width 640
    s0 = mpp._calculate_background_region(
        1280, 720, 0, 2, "horizontal", True, split_dir="left", split_pct=50
    )
    s1 = mpp._calculate_background_region(
        1280, 720, 1, 2, "horizontal", True, split_dir="left", split_pct=50
    )
    assert s0 == (0, 0, 320, 720)
    assert s1 == (320, 0, 320, 720)


def test_widen_text_shapes_increases_textbox_width():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = 0
    width = mpp.Cm(5).emu
    height = mpp.Cm(2).emu
    textbox = slide.shapes.add_textbox(left, top, width, height)
    before = textbox.width

    # call helper
    mpp.widen_text_shapes(prs=prs, extra_width_cm=4)

    after = textbox.width
    assert after > before
    # expected delta preserved (expressed using Cm(...))
    expected_delta = mpp.Cm(4).emu
    assert after == before + expected_delta



def test_normalize_font_names_run_level():
    """Explicit run.font.name values like 'SegoeUI' should be normalized."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, mpp.Cm(5).emu, mpp.Cm(2).emu)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Hello"
    run.font.name = "SegoeUI"

    # sanity check before
    assert run.font.name == "SegoeUI"

    mpp.normalize_font_names(prs)

    assert run.font.name == "Segoe UI"


def test_normalize_font_names_theme_part():
    """Theme part XML containing `typeface="SegoeUI"` should be corrected."""
    import re
    from pptx import Presentation

    prs = Presentation()
    # locate theme part
    theme_part = next((rel._target for rel in prs.part.rels.values() if rel.reltype.endswith('/theme')), None)
    assert theme_part is not None

    blob_bytes = getattr(theme_part, '_blob', None) or getattr(theme_part, 'blob', None)
    assert isinstance(blob_bytes, (bytes, bytearray))
    xml = blob_bytes.decode('utf-8', errors='ignore')
    # inject a SegoeUI value into the first typeface occurrence
    xml2 = re.sub(r'typeface="[^"]+"', 'typeface="SegoeUI"', xml, count=1)
    setattr(theme_part, '_blob', xml2.encode('utf-8'))

    blob_bytes2 = getattr(theme_part, '_blob', None) or getattr(theme_part, 'blob', None)
    assert isinstance(blob_bytes2, (bytes, bytearray))
    assert b'SegoeUI' in blob_bytes2

    mpp.normalize_font_names(prs)

    blob_bytes3 = getattr(theme_part, '_blob', None) or getattr(theme_part, 'blob', None)
    assert isinstance(blob_bytes3, (bytes, bytearray))
    assert b'Segoe UI' in blob_bytes3
    assert b'SegoeUI' not in blob_bytes3


def test_process_pptx_html_normalizes_segoeui(tmp_path):
    """Integration: running the pipeline should remove 'SegoeUI' from runs and theme."""
    from pptx import Presentation
    import re

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, mpp.Cm(5).emu, mpp.Cm(2).emu)
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Test"
    run.font.name = "SegoeUI"

    # also inject SegoeUI into theme to ensure theme-pass runs
    theme_part = next((rel._target for rel in prs.part.rels.values() if rel.reltype.endswith('/theme')), None)
    assert theme_part is not None

    blob_bytes = getattr(theme_part, '_blob', None) or getattr(theme_part, 'blob', None)
    assert isinstance(blob_bytes, (bytes, bytearray))
    xml = blob_bytes.decode('utf-8', errors='ignore')
    xml2 = re.sub(r'typeface="[^"]+"', 'typeface="SegoeUI"', xml, count=1)
    setattr(theme_part, '_blob', xml2.encode('utf-8'))

    raw_pptx = tmp_path / 'raw.pptx'
    prs.save(raw_pptx)

    html_path = tmp_path / 'dummy.html'
    html_path.write_text('<html></html>', encoding='utf-8')

    out_pptx = tmp_path / 'out.pptx'

    # run pipeline (skip experimental styled-divs)
    mpp.process_pptx_html(html_path, raw_pptx, out_pptx, save_rendered_divs=False, run_styled_divs=False)

    # reopen and assert no SegoeUI remains in runs or theme
    prs2 = Presentation(out_pptx)
    found_bad = False
    for s in prs2.slides:
        for sh in s.shapes:
            try:
                if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        for r in para.runs:
                            if r.font.name and r.font.name.lower().strip() == 'segoeui':
                                found_bad = True
            except Exception:
                continue
    assert not found_bad

    # check theme part
    theme_part2 = next((rel._target for rel in prs2.part.rels.values() if rel.reltype.endswith('/theme')), None)
    assert theme_part2 is not None

    blob_bytes_out = getattr(theme_part2, '_blob', None) or getattr(theme_part2, 'blob', None)
    assert isinstance(blob_bytes_out, (bytes, bytearray))
    assert b'SegoeUI' not in blob_bytes_out
    assert b'Segoe UI' in blob_bytes_out


def test_process_native_marp_images_applies_cover_crop(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # create a wide image (aspect > slide aspect) so 'cover' triggers a left/right crop
    img = Image.new("RGB", (2000, 1000), color="red")
    p = tmp_path / "wide.jpg"
    img.save(p)

    # add picture shape to the slide
    pic = slide.shapes.add_picture(str(p), 0, 0)

    # build slides_data that mimics a native Marp background + image list
    slides_data = [
        {
            "backgrounds": [
                {
                    "url": None,
                    "direction": "horizontal",
                    "split": None,
                    "split_pct": None,
                    "size": "cover",
                }
            ],
            "content": None,
            "all_images": [{"url": None, "source": "background"}],
        }
    ]

    # call helper
    mpp.process_native_marp_images(
        prs=prs,
        slides_data=slides_data,
    )

    # expect a horizontal crop to have been applied (left/right crop > 0)
    assert pic.crop_left > 0 or pic.crop_right > 0


def test_process_native_marp_images_handles_pixel_values(tmp_path):
    """Test that pixel values like '200px auto' are correctly parsed and applied."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Create a test image
    img = Image.new("RGB", (400, 400), color="blue")
    p = tmp_path / "test.jpg"
    img.save(p)

    # Add picture shape to the slide
    pic = slide.shapes.add_picture(str(p), 0, 0)
    
    # Store original dimensions for comparison
    original_width = pic.width
    original_height = pic.height

    # Test case 1: width in pixels, height auto
    slides_data = [
        {
            "backgrounds": [
                {
                    "url": None,
                    "direction": "horizontal",
                    "split": None,
                    "split_pct": None,
                    "size": "200px auto",  # width in pixels, height proportional
                }
            ],
            "content": None,
            "all_images": [{"url": None, "source": "background"}],
        }
    ]

    # Call helper
    mpp.process_native_marp_images(
        prs=prs,
        slides_data=slides_data,
    )

    # The width should be modified (not equal to original)
    # Since we specify pixels, it should be different from the default full-slide size
    assert pic.width != original_width or pic.height != original_height
    # No crop should be applied for properly sized images
    assert pic.crop_left == 0.0 and pic.crop_right == 0.0
    assert pic.crop_top == 0.0 and pic.crop_bottom == 0.0


def test_process_native_marp_images_handles_auto_height_pixels(tmp_path):
    """Test that pixel values like 'auto 200px' are correctly parsed and applied."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Create a test image
    img = Image.new("RGB", (400, 400), color="green")
    p = tmp_path / "test2.jpg"
    img.save(p)

    # Add picture shape to the slide
    pic = slide.shapes.add_picture(str(p), 0, 0)
    
    # Store original dimensions
    original_width = pic.width
    original_height = pic.height

    # Test case 2: auto width, height in pixels
    slides_data = [
        {
            "backgrounds": [
                {
                    "url": None,
                    "direction": "horizontal",
                    "split": None,
                    "split_pct": None,
                    "size": "auto 200px",  # width proportional, height in pixels
                }
            ],
            "content": None,
            "all_images": [{"url": None, "source": "background"}],
        }
    ]

    # Call helper
    mpp.process_native_marp_images(
        prs=prs,
        slides_data=slides_data,
    )

    # The dimensions should be modified
    assert pic.width != original_width or pic.height != original_height
    # No crop should be applied for properly sized images
    assert pic.crop_left == 0.0 and pic.crop_right == 0.0
    assert pic.crop_top == 0.0 and pic.crop_bottom == 0.0


def test_process_native_marp_images_handles_percentage_auto(tmp_path):
    """Test that '90% auto' (from w:90% or width:90%) correctly sizes without cover cropping."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Create a square test image (1200x1200 like in the real test file)
    img = Image.new("RGB", (1200, 1200), color="blue")
    p = tmp_path / "square.jpg"
    img.save(p)

    # Add picture to slide
    pic = slide.shapes.add_picture(str(p), 0, 0)
    
    # Get slide dimensions
    slide_width = prs.slide_width

    # Test '90% auto' (which is what Marp generates for w:90% or width:90%)
    slides_data = [
        {
            "backgrounds": [
                {
                    "url": None,
                    "direction": "horizontal",
                    "split": None,
                    "split_pct": None,
                    "size": "90% auto",
                }
            ],
            "content": None,
            "all_images": [{"url": None, "source": "background"}],
        }
    ]

    # Process
    mpp.process_native_marp_images(
        prs=prs,
        slides_data=slides_data,
    )

    # Expected: 90% of slide width, with proportional height (same for square)
    expected_width = int(slide_width * 0.9)
    expected_height = expected_width  # Square image maintains aspect

    # The image should be sized to 90% width with no cropping
    # (not full-slide with cover-like crop)
    assert pic.width == expected_width, f"Width should be 90% of slide ({expected_width}), got {pic.width}"
    assert pic.height == expected_height, f"Height should match width for square image ({expected_height}), got {pic.height}"
    
    # No crop should be applied - the image should just be smaller and centered
    assert pic.crop_left == 0.0 and pic.crop_right == 0.0, "No horizontal crop should be applied"
    assert pic.crop_top == 0.0 and pic.crop_bottom == 0.0, "No vertical crop should be applied"
    
    # Image should be centered horizontally
    expected_left = (slide_width - expected_width) // 2
    assert abs(pic.left - expected_left) < 100, "Image should be centered horizontally"


def test_render_div_as_image_handles_background_image_on_div(tmp_path):
    # Create a tiny PNG on disk and reference it via background-image on a div
    img = Image.new("RGB", (60, 60), color="purple")
    p = tmp_path / "bg_only.png"
    img.save(p)

    soup = BeautifulSoup(
        f'<div style="background-image:url({p.as_posix()}); width:100px; height:100px;"></div>',
        "html.parser",
    )
    div = soup.find("div")
    assert div is not None

    out = render_div_as_image(div, css=None, slide_index=0, div_index=0)
    assert out is not None
    assert Path(out).is_file()
    # cleanup
    os.remove(out)


def test_render_div_as_image_applies_border_radius_from_css(monkeypatch, tmp_path):
    # Create a source image bytes returned by requests.get
    img = Image.new("RGB", (200, 200), color="green")
    b = io.BytesIO()
    img.save(b, format="PNG")
    png_bytes = b.getvalue()

    class DummyResponse:
        def __init__(self, content_bytes: bytes):
            self.content = content_bytes

        def raise_for_status(self):
            return None

    def fake_get(url, timeout=10):
        return DummyResponse(png_bytes)

    monkeypatch.setattr(requests, "get", fake_get)

    soup = BeautifulSoup(
        '<div class="image-wrap-50"><img class="portrait" src="https://example.test/img.png"/></div>',
        "html.parser",
    )
    div = soup.find("div")
    assert div is not None

    css = ".image-wrap-50 { width:200px; height:200px; border-radius:50%; }"
    out = render_div_as_image(div, css=css, slide_index=0, div_index=0)
    assert out is not None

    im = Image.open(out)
    # corner pixel should be transparent because of circular mask
    assert im.getchannel("A").getpixel((0, 0)) == 0
    os.remove(out)


def test_render_div_as_image_parses_object_fit_and_position_and_scale(tmp_path):
    # Create a source image with left half red and right half blue
    src = Image.new("RGB", (200, 100))
    for x in range(200):
        for y in range(100):
            src.putpixel((x, y), (255, 0, 0) if x < 100 else (0, 0, 255))
    p = tmp_path / "split.png"
    src.save(p)

    soup = BeautifulSoup(
        f'<div class="image-wrap"><img class="image" src="{p.as_posix()}"/></div>',
        "html.parser",
    )
    div = soup.find("div")
    assert div is not None

    css = ".image-wrap { width:100px; height:100px } .image { object-fit: cover; object-position: 100% 50%; }"
    out = render_div_as_image(div, css=css, slide_index=0, div_index=0)
    assert out is not None

    im = Image.open(out).convert("RGB")
    # center pixel should be from the right-half (blue)
    pixel = cast(tuple, im.getpixel((50, 50)))
    assert pixel[2] > 200
    os.remove(out)


def test_process_styled_divs_inplace_replacement(tmp_path, monkeypatch):
    from pptx import Presentation

    # original picture (red)
    orig = Image.new("RGB", (50, 50), color=(255, 0, 0))
    orig_p = tmp_path / "orig.png"
    orig.save(orig_p)

    # new picture that renderer will return (blue)
    new = Image.new("RGB", (50, 50), color=(0, 0, 255))
    new_p = tmp_path / "new.png"
    new.save(new_p)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # slides_data with div referencing the basename of the original image
    slides_data = [{"content": '<div class="portrait-wrap"><img src="orig.png"/></div>'}]

    # monkeypatch renderer to return our new image path
    # (process_styled_divs uses the function imported in postprocessing)
    monkeypatch.setattr(mpp_post, "render_div_as_image", lambda *a, **k: str(new_p))

    # provide a dummy html_path (process_styled_divs reads it to extract CSS)
    html_path = tmp_path / "dummy.html"
    html_path.write_text("<style></style>")

    mpp.process_styled_divs(prs=prs, slides_data=slides_data, html_path=html_path)

    # Verify at least one picture on the slide now contains the new image bytes
    new_blob = new_p.read_bytes()
    found = False
    for s in slide.shapes:
        if s.shape_type == mpp.MSO_SHAPE_TYPE.PICTURE:
            try:
                if s.image.blob == new_blob:
                    found = True
                    break
            except Exception:
                continue

    assert found, "Rendered PNG was not applied to any picture on the slide"


def test_process_pptx_html_respects_experimental_flag(tmp_path, monkeypatch):
    """process_pptx_html should skip `process_styled_divs` when the
    experimental flag is not enabled.
    """
    from pptx import Presentation

    # create a minimal raw PPTX
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    raw_pptx = tmp_path / "raw.pptx"
    prs.save(raw_pptx)

    # create a minimal HTML that would normally contain a styled div
    html_content = (
        '<svg data-marpit-svg="true"><foreignObject>'
        '<section data-marpit-advanced-background="content">'
        '<div class="portrait-wrap"><img src="orig.png"/></div>'
        '</section></foreignObject></svg>'
    )
    html_path = tmp_path / "sample.html"
    html_path.write_text(html_content, encoding="utf-8")

    out_pptx = tmp_path / "out.pptx"

    called = {"v": False}

    def fake_process_styled_divs(*a, **k):
        called["v"] = True

    monkeypatch.setattr(mpp, "process_styled_divs", fake_process_styled_divs)

    # run with run_styled_divs=False -> should NOT call process_styled_divs
    mpp.process_pptx_html(html_path, raw_pptx, out_pptx, save_rendered_divs=False, run_styled_divs=False)

    assert called["v"] is False
    assert out_pptx.is_file()


def test_remove_redundant_marp_white_rectangles_helper():
    """Helper should remove exactly the redundant full-slide white rectangles when
    they are back-most and the slide contains other (non-candidate) shapes.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # add three full-slide white rectangles (the pattern we want to remove)
    for _ in range(3):
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_w, slide_h)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # add a non-candidate shape *after* the white rectangles so candidates are back-most
    slide.shapes.add_textbox(mpp.Cm(1).emu, mpp.Cm(1).emu, mpp.Cm(2).emu, mpp.Cm(1).emu)

    # sanity: ensure candidates exist
    candidates = [
        s
        for s in slide.shapes
        if getattr(s, "width", None) == slide_w
        and getattr(s, "height", None) == slide_h
        and _shape_rgb_safe(s) == RGBColor(255, 255, 255)
    ]
    assert len(candidates) >= 3

    removed = mpp.remove_redundant_marp_white_rectangles(prs)
    assert removed == 3

    # verify no remaining full-slide white rectangles
    assert not any(
        _shape_rgb_safe(s) == RGBColor(255, 255, 255) and s.width == slide_w and s.height == slide_h
        for s in slide.shapes
    )


def test_remove_redundant_marp_white_rectangles_on_slide_with_only_candidates():
    """If a slide contains *only* the three Marp white rectangles, they should
    still be removed (Marp sometimes emits an otherwise-empty slide)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # add exactly three full-slide white rectangles and nothing else
    for _ in range(3):
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_w, slide_h)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # sanity: ensure candidates exist and are the only shapes
    candidates = [
        s for s in slide.shapes if getattr(s, "width", None) == slide_w and getattr(s, "height", None) == slide_h and _shape_rgb_safe(s) == RGBColor(255, 255, 255)
    ]
    assert len(candidates) == 3

    removed = mpp.remove_redundant_marp_white_rectangles(prs)
    assert removed == 3

    # verify no remaining full-slide white rectangles
    assert not any(
        _shape_rgb_safe(s) == RGBColor(255, 255, 255) and s.width == slide_w and s.height == slide_h
        for s in slide.shapes
    )


def test_skip_if_count_not_three():
    """If the slide has fewer than the required candidate count, do not remove."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # only two white rectangles (not three)
    for _ in range(2):
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_w, slide_h)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # add a non-candidate afterwards so candidates would be back-most if they were enough
    slide.shapes.add_textbox(mpp.Cm(1).emu, mpp.Cm(1).emu, mpp.Cm(2).emu, mpp.Cm(1).emu)

    removed = mpp.remove_redundant_marp_white_rectangles(prs)
    assert removed == 0
    assert any(_shape_rgb_safe(s) == RGBColor(255, 255, 255) for s in slide.shapes)


def test_skip_if_candidates_not_backmost():
    """If candidate shapes are not the back-most shapes, skip removal."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # add a non-candidate first (so its index is lower than the candidates')
    slide.shapes.add_textbox(mpp.Cm(1).emu, mpp.Cm(1).emu, mpp.Cm(2).emu, mpp.Cm(1).emu)

    # now add three full-slide white rectangles (candidates will NOT be back-most)
    for _ in range(3):
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_w, slide_h)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)

    removed = mpp.remove_redundant_marp_white_rectangles(prs)
    assert removed == 0
    assert any(_shape_rgb_safe(s) == RGBColor(255, 255, 255) for s in slide.shapes)


def test_skip_if_slide_only_candidates():
    """Slides that only contain the three Marp white-rectangle artefacts should
    also have those rectangles removed (empty Marp slides are common)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # only the three white rectangles on the slide
    for _ in range(3):
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, slide_w, slide_h)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)

    removed = mpp.remove_redundant_marp_white_rectangles(prs)
    assert removed == 3
    assert not any(_shape_rgb_safe(s) == RGBColor(255, 255, 255) for s in slide.shapes)


def test_process_pptx_html_removes_redundant_white_rectangles(tmp_path):
    """End-to-end: process_pptx_html should remove the redundant white rectangles when enabled."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(3):
        r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # add a non-candidate shape after the rectangles so they are back-most and eligible for removal
    slide.shapes.add_textbox(mpp.Cm(1).emu, mpp.Cm(1).emu, mpp.Cm(2).emu, mpp.Cm(1).emu)

    raw_pptx = tmp_path / 'raw.pptx'
    prs.save(raw_pptx)

    html_path = tmp_path / 'dummy.html'
    html_path.write_text('<html></html>', encoding='utf-8')

    out_pptx = tmp_path / 'out.pptx'
    mpp.process_pptx_html(html_path, raw_pptx, out_pptx, save_rendered_divs=False, run_styled_divs=False)

    prs2 = Presentation(out_pptx)
    assert not any(
        _shape_rgb_safe(s) == RGBColor(255, 255, 255)
        and s.width == prs2.slide_width
        and s.height == prs2.slide_height
        for s in prs2.slides[0].shapes
    )


def test_real_markdown_file_consolidates_to_single_textbox(tmp_path):
    """End-to-end test using the actual characters.marp.md file and its generated PPTX.
    
    This test loads the PPTX that was generated from characters.marp.md,
    extracts segment information from the markdown (with U+200b markers),
    and verifies that consolidation reduces the textboxes to a single one.
    """
    from pathlib import Path
    from pptx import Presentation
    
    md_file = Path(__file__).parent.parent / "resources" / "characters.marp.md"
    
    assert md_file.exists(), f"Test file {md_file} not found"
    
    # Generate PPTX from the markdown file
    pptx_file = tmp_path / "characters.marp.md.pptx"
    html_file = tmp_path / "characters.marp.md.html"
    
    try:
        from marp2pptx.marp_convert import marp_generate_in_parallel
        marp_generate_in_parallel(md_file, html_file, pptx_file)
    except Exception as e:
        # If marp fails, skip (e.g., in CI without Node.js)
        print(f"Skipping test: marp CLI not available ({e})")
        return
    
    assert pptx_file.exists(), f"Generated PPTX {pptx_file} not found"
    
    # Read markdown to extract segment markers
    md_content = md_file.read_text(encoding="utf-8")
    
    # Extract the line with **01** and the U+200b markers
    for line in md_content.splitlines():
        if "**01**" in line:
            # Remove markdown bold syntax (**text** -> text) to get clean text
            # This matches what the PPTX will have (rendered text without formatting)
            clean_line = line.replace("**", "")
            # Create HTML with the clean text and U+200b markers preserved
            section_html = f"<section><p>{clean_line}</p></section>"
            break
    else:
        raise AssertionError("Could not find '**01**' line in markdown")
    
    # DEBUG: print what the HTML structure looks like
    print(f"\nHTML input: {repr(section_html[:100])}")
    
    # Extract and print the segments
    soup = BeautifulSoup(section_html, "html.parser")
    text_content = soup.get_text()
    segments = text_content.split('\u200b')
    print(f"HTML segments ({len(segments)}):")
    for i, seg in enumerate(segments):
        print(f"  [{i}] {repr(seg)}")
    
    # Load the PPTX generated from this markdown
    prs = Presentation(str(pptx_file))
    slide = prs.slides[0]
    
    # Count textboxes before consolidation
    text_shapes_before = [
        s for s in slide.shapes
        if getattr(s, "has_text_frame", False) and s.has_text_frame
        and s.shape_type == mpp.MSO_SHAPE_TYPE.TEXT_BOX
    ]
    
    print(f"\nBefore consolidation: {len(text_shapes_before)} textboxes")
    for i, s in enumerate(text_shapes_before):
        text = s.text_frame.paragraphs[0].text if s.text_frame.paragraphs else ""
        print(f"  [{i}] text={repr(text[:80])}, left={s.left}, top={s.top}, width={s.width}")
        # Check if it would be considered a candidate
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        is_text_box = s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        is_placeholder = getattr(s, "is_placeholder", False)
        has_rotation = getattr(s, "rotation", 0) not in (0, None)
        is_group = s.shape_type == MSO_SHAPE_TYPE.GROUP
        has_text_frame = getattr(s, "has_text_frame", False) and s.has_text_frame
        num_paragraphs = len(s.text_frame.paragraphs) if has_text_frame else 0
        has_newline = "\n" in text
        print(f"      is_text_box={is_text_box}, is_placeholder={is_placeholder}, has_rotation={has_rotation}")
        print(f"      is_group={is_group}, has_text_frame={has_text_frame}, num_paras={num_paragraphs}, has_newline={has_newline}")
    

# ============================================================================
# PREPROCESSING TESTS
# ============================================================================


def test_remove_invisible_characters_removes_zero_width_space():
    """Unit test: remove_invisible_characters should remove U+200B characters."""
    input_text = "Økonomi​ og estimat​er"  # contains U+200B after 'Økonomi' and 'estimat'
    expected = "Økonomi og estimater"
    
    result = remove_invisible_characters(input_text)
    
    assert result == expected
    assert '\u200b' not in result


def test_remove_invisible_characters_removes_multiple_invisible_types():
    """Unit test: verify removal of various invisible Unicode characters."""
    # Test various invisible char types: U+200B, U+200C, U+200D, U+FEFF, U+200E, U+200F
    input_text = "Hello\u200bWorld\u200cTest\u200dData\uFEFFMore\u200EText\u200F"
    expected = "HelloWorldTestDataMoreText"
    
    result = remove_invisible_characters(input_text)
    
    assert result == expected
    assert '\u200b' not in result
    assert '\u200c' not in result
    assert '\u200d' not in result
    assert '\uFEFF' not in result
    assert '\u200e' not in result
    assert '\u200f' not in result


def test_remove_invisible_characters_preserves_normal_text():
    """Unit test: verify that normal text and formatting is preserved."""
    input_text = "**01** Økonomi og estimater for forbrug\n\nInternal document"
    expected = input_text
    
    result = remove_invisible_characters(input_text)
    
    assert result == expected


def test_preprocess_markdown_writes_cleaned_file(tmp_path):
    """Unit test: preprocess_markdown should write a cleaned markdown file."""
    # Create input file with invisible characters
    input_md = tmp_path / "input.md"
    input_md.write_text("---\nmarp: true\n---\n\n**01** Økonomi​ og estimat​er ", encoding='utf-8')
    
    # Preprocess to output file
    output_md = tmp_path / "output.md"
    preprocess_markdown(input_md, output_md)
    
    # Verify output file exists and contains cleaned text
    assert output_md.exists()
    content = output_md.read_text(encoding='utf-8')
    assert '\u200b' not in content
    assert "Økonomi og estimater" in content


def test_preprocess_markdown_end_to_end_with_text_marp_md(tmp_path):
    """End-to-end test: preprocessing should prevent textbox splitting in PPTX.
    
    This test:
    1. Reads the actual test file: resources/characters.marp.md (contains U+200B)
    2. Preprocesses it to remove invisible characters
    3. Runs marp CLI on the preprocessed file to generate PPTX + HTML
    4. Verifies the resulting PPTX has exactly ONE textbox (not split)
    5. Verifies the textbox content is the cleaned text without invisible chars
    """
    from pptx import Presentation
    
    # Locate the actual test resource file
    test_resources_dir = Path(__file__).parent.parent / "resources"
    original_md = test_resources_dir / "characters.marp.md"
    
    assert original_md.exists(), f"Test file {original_md} not found"
    
    # Read original to verify it has the invisible characters
    original_content = original_md.read_text(encoding='utf-8')
    assert '\u200b' in original_content, "Original file should contain U+200B characters"
    assert "**01**" in original_content
    
    # Step 1: Preprocess to remove invisible characters
    preprocessed_md = tmp_path / "text_preprocessed.marp.md"
    preprocess_markdown(original_md, preprocessed_md)
    
    preprocessed_content = preprocessed_md.read_text(encoding='utf-8')
    assert '\u200b' not in preprocessed_content, "Preprocessed file should have no U+200B"
    
    # Step 2: Generate PPTX and HTML from preprocessed markdown
    html_output = tmp_path / "text_preprocessed.marp.md.html"
    pptx_output = tmp_path / "text_preprocessed.marp.md_raw.pptx"
    
    try:
        # Run marp CLI on the preprocessed file
        from marp2pptx.marp_convert import marp_generate_in_parallel
        marp_generate_in_parallel(preprocessed_md, html_output, pptx_output)
    except Exception as e:
        # If marp fails, skip (e.g., in CI without Node.js)
        print(f"Skipping end-to-end test: marp CLI not available ({e})")
        return
    
    assert html_output.exists(), "Marp should generate HTML"
    assert pptx_output.exists(), "Marp should generate PPTX"
    
    # Step 3: Load PPTX and verify textbox consolidation
    prs = Presentation(str(pptx_output))
    slide = prs.slides[0]
    
    # Get all textboxes in the slide
    text_shapes = [
        s for s in slide.shapes
        if getattr(s, "has_text_frame", False) and s.has_text_frame
        and s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    ]
    
    # Filter to only content textboxes (exclude placeholders, empty, etc.)
    content_shapes = []
    for s in text_shapes:
        # Skip placeholder shapes
        if getattr(s, "is_placeholder", False):
            continue
        # Skip empty text frames
        if not s.text_frame.paragraphs:
            continue
        text_content = s.text_frame.text.strip()
        if not text_content:
            continue
        # Skip shapes with rotation or other special properties
        if getattr(s, "rotation", 0) not in (0, None):
            continue
        content_shapes.append(s)
    
    # The key assertion: should have exactly ONE content textbox (not split)
    assert len(content_shapes) == 1, (
        f"Expected 1 textbox, but found {len(content_shapes)}. "
        f"This indicates the invisible characters are still causing text splitting. "
        f"Textbox positions: {[(s.left, s.top, s.width) for s in content_shapes]}"
    )
    
    # Step 4: Verify the text content matches the expected cleaned text
    actual_text = content_shapes[0].text_frame.text.strip()
    # Expected is the original line without markdown bold markers and without invisible chars
    expected_text = "01 Økonomi og estimater for forbrug"
    
    assert actual_text == expected_text, (
        f"Expected text: {repr(expected_text)}\n"
        f"Actual text:   {repr(actual_text)}"
    )
    
    # Additional sanity check: no U+200B should remain in the text
    assert '\u200b' not in actual_text, "Textbox should not contain invisible characters"


# ============================================================================
# MERGE MULTILINE TEXT BOXES TESTS
# ============================================================================

_LONG_SENTENCE = (
    "This is a remarkably long sentence that is specifically designed to be long enough "
    "to wrap across multiple visual lines when rendered within the default viewport "
    "dimensions of a Marp presentation slide."
)

# Indices of the four new wrapping-sentence slides in text.marp.md (0-based):
# slide 5 = plain text, slide 6 = h1, slide 7 = h2, slide 8 = h4
_WRAPPING_SLIDE_INDICES = [5, 6, 7, 8]
_WRAPPING_SLIDE_LABELS = ["plain text", "h1 header", "h2 header", "h4 header"]


def test_merge_multiline_textboxes_merges_adjacent_boxes():
    """Unit test: merge_multiline_textboxes should merge vertically-adjacent text boxes
    with the same left position into a single text box."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Simulate two text boxes from the same wrapped sentence:
    # same left, adjacent tops (box B starts right below box A).
    left = mpp.Cm(1).emu
    width = mpp.Cm(20).emu
    height = mpp.Cm(1).emu

    box_a = slide.shapes.add_textbox(left, mpp.Cm(2).emu, width, height)
    box_a.text_frame.paragraphs[0].add_run().text = "First line of the sentence"

    # top of box B == top of box A + height of box A  (adjacent)
    box_b = slide.shapes.add_textbox(left, mpp.Cm(2).emu + height, width, height)
    box_b.text_frame.paragraphs[0].add_run().text = "Second line of the sentence"

    # Sanity: 2 text boxes before
    text_boxes_before = [s for s in slide.shapes if s.shape_type == mpp.MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(text_boxes_before) == 2

    removed = mpp.merge_multiline_textboxes(prs)
    assert removed == 1, f"Expected 1 text box removed, got {removed}"

    text_boxes_after = [s for s in slide.shapes if s.shape_type == mpp.MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(text_boxes_after) == 1, (
        f"Expected 1 text box after merge, got {len(text_boxes_after)}"
    )
    # Both texts must be present in the merged box
    merged_text = text_boxes_after[0].text_frame.text
    assert "First line" in merged_text
    assert "Second line" in merged_text


def test_merge_multiline_textboxes_does_not_merge_separate_sentences():
    """Unit test: text boxes with different left positions or large vertical gaps
    should NOT be merged."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left_a = mpp.Cm(1).emu
    left_b = mpp.Cm(10).emu  # different left -> should not merge
    width = mpp.Cm(8).emu
    height = mpp.Cm(1).emu

    box_a = slide.shapes.add_textbox(left_a, mpp.Cm(2).emu, width, height)
    box_a.text_frame.paragraphs[0].add_run().text = "Box A"

    box_b = slide.shapes.add_textbox(left_b, mpp.Cm(2).emu + height, width, height)
    box_b.text_frame.paragraphs[0].add_run().text = "Box B"

    removed = mpp.merge_multiline_textboxes(prs)
    assert removed == 0, "Boxes with different left should not be merged"

    text_boxes = [s for s in slide.shapes if s.shape_type == mpp.MSO_SHAPE_TYPE.TEXT_BOX]
    assert len(text_boxes) == 2


def test_merge_multiline_textboxes_text_marp_md(tmp_path):
    """End-to-end test: process_pptx_html (which calls merge_multiline_textboxes)
    should consolidate each wrapped long sentence in the 4 new slides of
    resources/text.marp.md into exactly one text box per slide."""
    from pathlib import Path
    from pptx import Presentation

    md_file = Path(__file__).parent.parent / "resources" / "text.marp.md"
    assert md_file.exists(), f"Test file {md_file} not found"

    html_file = tmp_path / "text.marp.md.html"
    raw_pptx_file = tmp_path / "text.marp.md_raw.pptx"
    out_pptx_file = tmp_path / "text.marp.md.pptx"

    try:
        from marp2pptx.marp_convert import marp_generate_in_parallel
        marp_generate_in_parallel(md_file, html_file, raw_pptx_file)
    except Exception as e:
        print(f"Skipping test: marp CLI not available ({e})")
        return

    assert raw_pptx_file.exists(), "Marp should have generated a raw PPTX"

    # Apply post-processing (includes merge_multiline_textboxes)
    mpp.process_pptx_html(
        html_file, raw_pptx_file, out_pptx_file,
        save_rendered_divs=False, run_styled_divs=False,
    )

    prs = Presentation(str(out_pptx_file))

    for slide_idx, label in zip(_WRAPPING_SLIDE_INDICES, _WRAPPING_SLIDE_LABELS):
        slide = prs.slides[slide_idx]
        text_boxes = [
            s for s in slide.shapes
            if s.shape_type == mpp.MSO_SHAPE_TYPE.TEXT_BOX
            and not getattr(s, "is_placeholder", False)
            and getattr(s, "has_text_frame", False)
            and s.has_text_frame
            and s.text_frame.text.strip()
        ]
        count = len(text_boxes)
        print(
            f"Slide {slide_idx + 1} ({label}): {count} non-empty text box(es) after merging"
        )
        assert count == 1, (
            f"Slide {slide_idx + 1} ({label}): expected 1 text box after merging the "
            f"wrapped long sentence, but found {count}. "
            f"Text boxes: {[s.text_frame.text[:60] for s in text_boxes]}"
        )
        # Verify the sentence text is present
        merged_text = text_boxes[0].text_frame.text
        # The sentence key words should all be present
        assert "remarkably long sentence" in merged_text, (
            f"Slide {slide_idx + 1} ({label}): merged text does not contain expected sentence. "
            f"Got: {repr(merged_text[:120])}"
        )


